# -*- coding: utf-8 -*-
import os
import sys
import csv
import time
import ctypes
import threading
import queue
import socket
import uuid
import getpass
from pathlib import Path
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
import customtkinter as ctk
from tkinter import filedialog, messagebox

# 文档解析与OCR库 (打包时需确保已安装)
try:
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation
    from rapidocr_onnxruntime import RapidOCR
except ImportError as e:
    print(
        f"缺少依赖库: {e}\n请执行: pip install python-docx openpyxl python-pptx pdfplumber pdf2image rapidocr_onnxruntime"
    )
    sys.exit(1)

# ==================== 系统信息获取 ====================
def get_system_info():
    """获取电脑用户名、电脑名、MAC地址"""
    username = getpass.getuser()
    hostname = socket.gethostname()
    mac_raw = '%012x' % uuid.getnode()
    mac_address = ':'.join(mac_raw[i:i+2] for i in range(0, 12, 2))
    return username, hostname, mac_address

# ==================== 路径与日志配置 ====================
if getattr(sys, 'frozen', False):
    BASE_DIR = os.path.dirname(sys.executable)
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

LOG_FILE = os.path.join(BASE_DIR, 'search_log.txt')

import logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
    handlers=[logging.FileHandler(LOG_FILE, encoding='utf-8')]
)
logger = logging.getLogger(__name__)

FILE_ATTRIBUTE_HIDDEN = 0x2
FILE_ATTRIBUTE_SYSTEM = 0x4

# ==================== 核心搜索引擎 ====================
class SearchEngine:
    # 所有支持的扩展名（用于GUI展示）
    ALL_SUPPORTED_EXTENSIONS = {
        '.docx', '.xlsx', '.pptx', '.txt', '.log', '.md', '.csv',
        '.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp', '.pdf'
    }

    SKIP_DIRS = {
        '$recycle.bin', '$windows.~bt', '$windows.~ws', 'windows',
        'program files', 'program files (x86)', 'programdata', 'recovery',
        'system volume information', 'msocache', 'intel', 'amd',
        '.git', 'node_modules', '__pycache__', '.vscode', '.idea', 'appdata'
    }

    def __init__(self, keywords, scan_path, output_path, max_workers, log_queue, selected_extensions=None):
        self.keywords = keywords
        self.scan_path = scan_path
        self.output_path = output_path
        self.max_workers = max_workers
        self.log_queue = log_queue
        # ✅ 新增：用户选择的文件类型，None表示全选
        self.selected_extensions = selected_extensions or self.ALL_SUPPORTED_EXTENSIONS
        self._ocr_engine = None
        self._poppler_path = self._find_poppler()
        self.stats = {'total': 0, 'scanned': 0, 'matched': 0, 'error': 0}
        self.results = []
        self.is_running = False
        # ✅ 获取系统信息（只获取一次）
        self.sys_username, self.sys_hostname, self.sys_mac = get_system_info()

    def _send_log(self, msg):
        self.log_queue.put(msg)

    @property
    def ocr_engine(self):
        if self._ocr_engine is None:
            self._send_log("⏳ 正在初始化 OCR 引擎 (首次较慢)...")
            self._ocr_engine = RapidOCR()
            self._send_log("✅ OCR 引擎就绪")
            self._send_log("🔍 正在扫描文件内容...")
        return self._ocr_engine

    def _find_poppler(self):
        for p in [os.path.join(BASE_DIR, 'poppler', 'bin'), os.path.join(os.getcwd(), 'poppler', 'bin')]:
            if os.path.exists(p):
                return p
        return None

    def _is_skip_dir(self, dir_path):
        name = os.path.basename(dir_path).lower()
        if name in self.SKIP_DIRS:
            return True
        if sys.platform == 'win32':
            try:
                attrs = ctypes.windll.kernel32.GetFileAttributesW(str(dir_path))
                if attrs != -1 and (attrs & (FILE_ATTRIBUTE_HIDDEN | FILE_ATTRIBUTE_SYSTEM)):
                    return True
            except Exception:
                pass
        return False

    # --- 文本提取方法 ---
    def extract_txt(self, fp):
        for enc in ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']:
            try:
                with open(fp, 'r', encoding=enc) as f:
                    return f.read()
            except Exception:
                continue
        return ""

    def extract_docx(self, fp):
        doc = Document(fp)
        texts = [p.text for p in doc.paragraphs if p.text.strip()]
        for t in doc.tables:
            for r in t.rows:
                for c in r.cells:
                    if c.text.strip():
                        texts.append(c.text)
        return '\n'.join(texts)

    def extract_xlsx(self, fp):
        try:
            wb = load_workbook(fp, read_only=True, data_only=True)
            texts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            texts.append(str(cell.value))
            wb.close()
            return '\n'.join(texts)
        except Exception:
            return ""

    def extract_pptx(self, fp):
        prs = Presentation(fp)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.extend([p.text for p in shape.text_frame.paragraphs if p.text.strip()])
                if shape.has_table:
                    for row in shape.table.rows:
                        texts.extend([c.text for c in row.cells if c.text.strip()])
        return '\n'.join(texts)

    def extract_image(self, fp):
        res, _ = self.ocr_engine(fp)
        return '\n'.join([l[1] for l in res]) if res else ""

    def extract_pdf(self, fp):
        import pdfplumber
        from pdf2image import convert_from_path
        parts = []
        try:
            with pdfplumber.open(fp) as pdf:
                for page in pdf.pages:
                    parts.append(page.extract_text() or "")
                    for tbl in (page.extract_tables() or []):
                        for row in tbl:
                            cells = [str(c).strip() for c in row if c]
                            if cells:
                                parts.append(' | '.join(cells))
        except Exception as e:
            self._send_log(f"⚠️ PDF文本提取失败: {e}")

        full = '\n'.join(parts).strip()
        clean = full.replace('\n', '').replace(' ', '')
        avg = len(clean) / max(1, len(parts))

        if avg < 10:
            self._send_log(f"📄 疑似扫描件，启用OCR: {os.path.basename(fp)}")
            try:
                imgs = convert_from_path(fp, dpi=200, poppler_path=self._poppler_path)
                ocr_parts = []
                for img in imgs:
                    res, _ = self.ocr_engine(img)
                    if res:
                        ocr_parts.append('\n'.join([l[1] for l in res]))
                if ocr_parts:
                    full = '\n'.join(ocr_parts)
            except Exception as e:
                self._send_log(f"❌ PDF OCR失败: {e}")
        return full

    def _check_keywords(self, text):
        if not text:
            return []
        lower_text = text.lower()
        return [kw for kw in self.keywords if kw.lower() in lower_text]

    def _get_context(self, text, kw, chars=50):
        idx = text.lower().find(kw.lower())
        if idx == -1:
            return ""
        s, e = max(0, idx - chars), min(len(text), idx + len(kw) + chars)
        ctx = text[s:e].replace('\n', ' ').replace('\r', ' ').strip()
        return f"{'...' if s > 0 else ''}{ctx}{'...' if e < len(text) else ''}"

    def process_file(self, fp):
        ext = Path(fp).suffix.lower()
        res = {'path': fp, 'keywords': [], 'contexts': {}, 'error': None}
        try:
            text = ""
            if ext in {'.txt', '.log', '.md', '.csv'}:
                text = self.extract_txt(fp)
            elif ext == '.docx':
                text = self.extract_docx(fp)
            elif ext == '.xlsx':
                text = self.extract_xlsx(fp)
            elif ext == '.pptx':
                text = self.extract_pptx(fp)
            elif ext in {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp'}:
                text = self.extract_image(fp)
            elif ext == '.pdf':
                text = self.extract_pdf(fp)

            matched = self._check_keywords(text)
            if matched:
                res['keywords'] = matched
                res['contexts'] = {k: self._get_context(text, k) for k in matched}
        except Exception as e:
            res['error'] = str(e)
        return res

    def run(self):
        self.is_running = True
        self.stats = {'total': 0, 'scanned': 0, 'matched': 0, 'error': 0}
        self.results = []
        start_time = datetime.now()

        self._send_log(f"🚀 开始扫描: {self.scan_path}")
        self._send_log(f"🔑 关键词: {', '.join(self.keywords)}")
        self._send_log(f"📋 扫描类型: {', '.join(sorted(self.selected_extensions))}")
        self._send_log(f"💻 用户: {self.sys_username} | 电脑: {self.sys_hostname} | MAC: {self.sys_mac}")

        files = []
        for root, dirs, fnames in os.walk(self.scan_path):
            dirs[:] = [d for d in dirs if not self._is_skip_dir(os.path.join(root, d))]
            for f in fnames:
                # ✅ 仅扫描用户选择的文件类型
                if Path(f).suffix.lower() in self.selected_extensions:
                    files.append(os.path.join(root, f))

        self.stats['total'] = len(files)
        self._send_log(f"📂 发现 {len(files)} 个目标文件")
        self.log_queue.put(("PROGRESS_MAX", len(files)))

        scan_start_time = time.time()

        with ThreadPoolExecutor(max_workers=self.max_workers) as ex:
            futures = {ex.submit(self.process_file, f): f for f in files}
            for future in as_completed(futures):
                if not self.is_running:
                    break
                r = future.result()
                self.stats['scanned'] += 1
                self.log_queue.put(("PROGRESS_VAL", self.stats['scanned']))

                if r['error']:
                    self.stats['error'] += 1
                    self._send_log(f"❌ 出错: {r['path']} - {r['error']}")
                elif r['keywords']:
                    self.stats['matched'] += 1
                    self.results.append(r)
                    self._send_log(f"✅ 命中: {Path(r['path']).name} -> {', '.join(r['keywords'])}")

                scanned = self.stats['scanned']
                total = self.stats['total']
                is_checkpoint = (scanned % 5 == 0) or (scanned == total)
                just_matched = bool(r['keywords']) and not r['error']

                if is_checkpoint or just_matched:
                    elapsed = time.time() - scan_start_time
                    speed = scanned / elapsed if elapsed > 0 else 0
                    eta = (total - scanned) / speed if speed > 0 else 0
                    self._send_log(
                        f" 📄 进度: {scanned}/{total} | "
                        f"速度: {speed:.1f}个/s | "
                        f"预计剩余: {eta:.0f}s"
                    )

        if self.results and self.is_running:
            # ✅ 新增：记录扫描完成时间（所有文件处理完毕后统一获取）
            scan_time = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            out_file = os.path.join(self.output_path, f"搜索结果_{datetime.now():%Y%m%d_%H%M%S}.csv")
            try:
                with open(out_file, 'w', newline='', encoding='utf-8-sig') as f:
                    w = csv.writer(f)
                    # ✅ 表头末尾新增「扫描时间」
                    w.writerow(['文件名', '文件类型', '文件路径', '匹配关键词', '上下文摘要',
                                '电脑用户名', '电脑名', 'MAC地址', '扫描时间'])
                    for r in self.results:
                        p = Path(r['path'])
                        file_name = p.name
                        file_type = p.suffix.lstrip('.')
                        kw_str = ', '.join(r['keywords'])
                        ctx_str = ' | '.join(
                            f"[{k}] {r['contexts'].get(k, '')}" for k in r['keywords']
                        )
                        # ✅ 每行末尾写入统一的扫描时间
                        w.writerow([file_name, file_type, r['path'], kw_str, ctx_str,
                                    self.sys_username, self.sys_hostname, self.sys_mac, scan_time])
                self._send_log(f"💾 结果已保存: {out_file}")
            except Exception as e:
                self._send_log(f"❌ 保存结果失败: {e}")

        duration = (datetime.now() - start_time).total_seconds()
        self._send_log(
            f"🏁 完成! 耗时:{duration:.1f}s | 总计:{self.stats['total']} | 命中:{self.stats['matched']} | 错误:{self.stats['error']}")
        self.log_queue.put(("DONE", None))


# ==================== GUI 界面 ====================
class SearchApp(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("智能文件关键词扫描工具")
        self.geometry("750x720")  # ✅ 增大窗口以容纳文件类型选择框
        self.minsize(600, 600)

        ctk.set_appearance_mode("System")
        ctk.set_default_color_theme("blue")

        self.engine = None
        self.log_queue = queue.Queue()
        self.worker_thread = None
        self.ext_vars = {}  # ✅ 存储各文件类型的勾选变量

        self._build_ui()
        self.after(100, self._poll_queue)

    def _build_ui(self):
        # ✅ 顶部功能说明区域
        info_frame = ctk.CTkFrame(self, fg_color="transparent")
        info_frame.pack(fill="x", padx=15, pady=(15, 5))

        info_text = (
            "📌 工具说明：本工具用于在指定文件夹中批量搜索包含关键词的文件内容。\n"
            "📥 输入：扫描目录、结果输出目录、搜索关键词（多个关键词用中文或英文逗号隔开）、并发线程数。\n"
            "📤 返回：匹配文件的路径、命中关键词、关键词前后50字上下文摘要，结果自动保存为CSV文件。\n"
            "📄 支持格式：\n"
            "   • 文档类：.docx / .xlsx / .pptx / .pdf（含扫描件OCR识别）\n"
            "   • 文本类：.txt / .log / .md / .csv\n"
            "   • 图片类：.jpg / .jpeg / .png / .bmp / .tiff / .webp（OCR识别）"
        )

        self.info_box = ctk.CTkTextbox(
            info_frame,
            height=130,
            font=("Microsoft YaHei UI", 12),
            fg_color=("gray90", "gray17"),
            text_color=("gray10", "gray90")
        )
        self.info_box.pack(fill="x")
        self.info_box.insert("1.0", info_text)
        self.info_box.configure(state="disabled")

        # --- 输入区域 ---
        frame_input = ctk.CTkFrame(self)
        frame_input.pack(fill="x", padx=15, pady=(5, 5))

        # 扫描路径
        row1 = ctk.CTkFrame(frame_input, fg_color="transparent")
        row1.pack(fill="x", pady=5)
        ctk.CTkLabel(row1, text="扫描路径:", width=80).pack(side="left")
        self.entry_scan = ctk.CTkEntry(row1)
        self.entry_scan.pack(side="left", fill="x", expand=True, padx=(0, 5))
        self.entry_scan.insert(0, r"D:\Documents")
        ctk.CTkButton(row1, text="浏览", width=60, command=self._browse_scan).pack(side="right")

        # 输出路径
        row2 = ctk.CTkFrame(frame_input, fg_color="transparent")
        row2.pack(fill="x", pady=5)
        ctk.CTkLabel(row2, text="输出路径:", width=80).pack(side="left")
        self.entry_output = ctk.CTkEntry(row2)
        self.entry_output.pack(side="left", fill="x", expand=True, padx=(0, 5))
        self.entry_output.insert(0, BASE_DIR)
        ctk.CTkButton(row2, text="浏览", width=60, command=self._browse_output).pack(side="right")

        # 关键词
        row3 = ctk.CTkFrame(frame_input, fg_color="transparent")
        row3.pack(fill="x", pady=5)
        ctk.CTkLabel(row3, text="关键词:", width=80).pack(side="left")
        self.entry_kw = ctk.CTkEntry(row3, placeholder_text="请输入搜索关键词")
        self.entry_kw.pack(side="left", fill="x", expand=True, padx=(0, 8))
        self.entry_kw.insert(0, "秘密,机密,绝密")
        ctk.CTkLabel(
            row3,
            text="(多个关键词用中文或英文逗号隔开)",
            text_color="gray",
            font=("Microsoft YaHei UI", 12)
        ).pack(side="right")

        # 线程数
        row4 = ctk.CTkFrame(frame_input, fg_color="transparent")
        row4.pack(fill="x", pady=5)
        ctk.CTkLabel(row4, text="线程数:", width=80).pack(side="left")
        self.entry_threads = ctk.CTkEntry(row4, width=80)
        self.entry_threads.insert(0, "4")
        self.entry_threads.pack(side="left")
        ctk.CTkLabel(row4, text="(推荐4-8，SSD固态硬盘可设更高)", text_color="gray").pack(side="left", padx=10)

        # ✅ 新增：文件类型选择框
        row5 = ctk.CTkFrame(frame_input, fg_color="transparent")
        row5.pack(fill="x", pady=5)
        ctk.CTkLabel(row5, text="扫描类型:", width=80).pack(side="left", anchor="n", pady=(5, 0))

        ext_frame = ctk.CTkFrame(row5, fg_color="transparent")
        ext_frame.pack(side="left", fill="x", expand=True)

        # 按类别分组排列
        categories = [
            ("文档", ['.docx', '.xlsx', '.pptx', '.pdf']),
            ("文本", ['.txt', '.log', '.md', '.csv']),
            ("图片", ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp']),
        ]

        for cat_name, exts in categories:
            cat_row = ctk.CTkFrame(ext_frame, fg_color="transparent")
            cat_row.pack(fill="x", pady=1)
            ctk.CTkLabel(cat_row, text=f"{cat_name}:", width=40, text_color="gray",
                         font=("Microsoft YaHei UI", 11)).pack(side="left")
            for ext in exts:
                var = ctk.BooleanVar(value=True)  # ✅ 默认全选
                self.ext_vars[ext] = var
                cb = ctk.CTkCheckBox(
                    cat_row, text=ext, variable=var,
                    font=("Microsoft YaHei UI", 11),
                    width=65, height=22
                )
                cb.pack(side="left", padx=2)

        # 全选/取消全选按钮
        btn_row = ctk.CTkFrame(ext_frame, fg_color="transparent")
        btn_row.pack(fill="x", pady=(3, 0))
        ctk.CTkButton(btn_row, text="全选", width=50, height=24,
                      font=("Microsoft YaHei UI", 11),
                      command=self._select_all_ext).pack(side="left", padx=(42, 3))
        ctk.CTkButton(btn_row, text="取消全选", width=60, height=24,
                      font=("Microsoft YaHei UI", 11),
                      command=self._deselect_all_ext).pack(side="left")

        # --- 按钮区 ---
        btn_frame = ctk.CTkFrame(self, fg_color="transparent")
        btn_frame.pack(fill="x", padx=15, pady=5)
        self.btn_start = ctk.CTkButton(btn_frame, text="▶ 开始搜索", height=36, command=self._start_search)
        self.btn_start.pack(side="left", expand=True, fill="x", padx=(0, 5))
        self.btn_stop = ctk.CTkButton(btn_frame, text="⏹ 停止", height=36, width=100,
                                      fg_color="red", hover_color="darkred",
                                      command=self._stop_search, state="disabled")
        self.btn_stop.pack(side="right")

        # --- 进度条 ---
        self.progress = ctk.CTkProgressBar(self)
        self.progress.pack(fill="x", padx=15, pady=(0, 5))
        self.progress.set(0)

        # --- 日志区 ---
        self.log_box = ctk.CTkTextbox(self, font=("Consolas", 12))
        self.log_box.pack(fill="both", expand=True, padx=15, pady=(0, 15))
        self.log_box.configure(state="disabled")

    # ✅ 全选/取消全选方法
    def _select_all_ext(self):
        for var in self.ext_vars.values():
            var.set(True)

    def _deselect_all_ext(self):
        for var in self.ext_vars.values():
            var.set(False)

    def _browse_scan(self):
        path = filedialog.askdirectory(title="选择扫描目录")
        if path:
            self.entry_scan.delete(0, "end")
            self.entry_scan.insert(0, path)

    def _browse_output(self):
        path = filedialog.askdirectory(title="选择结果输出目录")
        if path:
            self.entry_output.delete(0, "end")
            self.entry_output.insert(0, path)

    def _append_log(self, msg):
        self.log_box.configure(state="normal")
        self.log_box.insert("end", msg + "\n")
        self.log_box.see("end")
        self.log_box.configure(state="disabled")

    def _poll_queue(self):
        try:
            while True:
                item = self.log_queue.get_nowait()
                if isinstance(item, tuple):
                    cmd, val = item
                    if cmd == "PROGRESS_MAX":
                        self.progress.set(0)
                    elif cmd == "PROGRESS_VAL" and self.engine and self.engine.stats['total'] > 0:
                        self.progress.set(val / self.engine.stats['total'])
                    elif cmd == "DONE":
                        self._on_done()
                else:
                    self._append_log(str(item))
        except queue.Empty:
            pass
        self.after(100, self._poll_queue)

    def _start_search(self):
        scan = self.entry_scan.get().strip()
        output = self.entry_output.get().strip()
        kw_raw = self.entry_kw.get().strip()
        threads = self.entry_threads.get().strip()

        if not os.path.isdir(scan):
            return messagebox.showerror("错误", "扫描路径不存在")
        if not os.path.isdir(output):
            return messagebox.showerror("错误", "输出路径不存在")
        if not kw_raw:
            return messagebox.showwarning("提示", "请输入至少一个关键词")

        keywords = [k.strip() for k in kw_raw.replace('，', ',').split(',') if k.strip()]
        if not keywords:
            return messagebox.showwarning("提示", "关键词格式不正确")

        # ✅ 获取用户选择的文件类型
        selected_exts = {ext for ext, var in self.ext_vars.items() if var.get()}
        if not selected_exts:
            return messagebox.showwarning("提示", "请至少选择一种文件类型")

        try:
            max_w = max(1, min(int(threads), 32))
        except ValueError:
            max_w = 4

        # 清空旧日志并重置UI
        self.log_box.configure(state="normal")
        self.log_box.delete("1.0", "end")
        self.log_box.configure(state="disabled")
        self.progress.set(0)
        self.btn_start.configure(state="disabled")
        self.btn_stop.configure(state="normal")

        # ✅ 传入选择的文件类型
        self.engine = SearchEngine(keywords, scan, output, max_w, self.log_queue, selected_exts)
        self.worker_thread = threading.Thread(target=self.engine.run, daemon=True)
        self.worker_thread.start()

    def _stop_search(self):
        if self.engine:
            self.engine.is_running = False
            self._append_log("⚠️ 用户手动停止搜索...")
            self.btn_stop.configure(state="disabled")

    def _on_done(self):
        self.btn_start.configure(state="normal")
        self.btn_stop.configure(state="disabled")
        self.progress.set(1)


if __name__ == "__main__":
    app = SearchApp()
    app.mainloop()